from flask import Flask, render_template, request, redirect, url_for, jsonify

import PyPDF2
from pptx import Presentation
from docx import Document
import textract
import google.generativeai as genai

app = Flask(_name_)

genai.configure(api_key="AIzaSyAqxVIKyLPplBCNp_QLXPF0T1gAmuh7wbk")

generation_config = {
    "temperature": 0.9,
    "top_p": 1,
    "top_k": 1,
    "max_output_tokens": 2048,
}

safety_settings = [
    {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
]

model = genai.GenerativeModel(
    model_name="gemini-1.0-pro",
    generation_config=generation_config,
    safety_settings=safety_settings,
)
initial_summary_generated = False
file_text = ""

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    global initial_summary_generated, file_text  # Uncomment this line

    if request.method == 'POST':
        try:
            print("Received POST request")
            
            # if 'file' not in request.files:
            #     print("No file part in the request")
            #     return jsonify({"error": "No file part"}), 400

            # if 'file' in request.files:
            file = request.files['file']


            if file.filename != '':
                
                if file.filename.endswith('.pdf'):
                    file_text = extract_text_from_pdf(file)
                elif file.filename.endswith(('.ppt', '.pptx')):
                    file_text = extract_text_from_ppt(file)
                elif file.filename.endswith('.txt'):
                    file_text = extract_text_from_txt(file)
                elif file.filename.endswith(('.doc', '.docx')):
                    file_text = extract_text_from_doc(file)  
                elif file.filename.endswith('.tex'):
                    file_text = extract_text_from_tex(file)        
                else:
                    return jsonify({"error": "Unsupported file format"}), 400
    

                if not initial_summary_generated:
                    print("Generating initial summary")
                    # Generate summary using the language model for the first time
                    convo = model.start_chat(history=[])
                    convo.send_message(file_text + "Summerise the above")
                    model_response = convo.last.text
                    initial_summary_generated = True
            print("Returning response")
            return jsonify({"extracted_text": file_text, "model_response": model_response})    

        except Exception as e:
            print(f"Exception: {str(e)}")  # Log the exception for debugging
            return jsonify({"error": str(e)}), 500

    # Handle GET request (render your HTML template)
    return render_template('index.html', user_question="", model_response="", extracted_text="")

@app.route('/userQues', methods=['POST'])
def chat():
    try:
        user_question = request.form['user_question']
        convo = model.start_chat(history=[])
        response = convo.send_message(file_text + " Use above text and give the following question answer: " + user_question)
        model_response = convo.last.text
        return jsonify({"model_response": model_response})
    except Exception as e:
        print(f"Exception: {str(e)}")
        return jsonify({"error": str(e)}), 500

def extract_text_from_pdf(file):
    text = ""
    with file.stream as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        for page_num in range(len(reader.pages)):
            text += reader.pages[page_num].extract_text()
    return text

def extract_text_from_ppt(file):
    text = ""
    presentation = Presentation(file.stream)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text

def extract_text_from_doc(file):
    text = ""
    doc = Document(file.stream)
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text

def extract_text_from_txt(file):
    text = ""
    with file.stream as txt_file:
        text = txt_file.read()
    return text    

def extract_text_from_tex(file):
    
    text = textract.process(file.stream, method='pdftotext')
    return text.decode('utf-8')
    

if _name_ == '_main_':
    app.run(debug=True)